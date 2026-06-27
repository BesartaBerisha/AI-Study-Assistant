import io

MAX_CHARS = 50_000


def parse_file(file_bytes: bytes, file_name: str) -> str:
    ext = file_name.lower().rsplit(".", 1)[-1]

    if ext == "pdf":
        return _parse_pdf(file_bytes)
    elif ext in ("ppt", "pptx"):
        return _parse_pptx(file_bytes)
    elif ext in ("doc", "docx"):
        return _parse_docx(file_bytes)
    elif ext in ("txt", "md"):
        return _parse_text(file_bytes)
    else:
        raise ValueError(f"Unsupported file type: .{ext}")


def _parse_pdf(file_bytes: bytes) -> str:
    import pypdf

    reader = pypdf.PdfReader(io.BytesIO(file_bytes))
    text = "\n\n".join(page.extract_text() or "" for page in reader.pages)
    return _truncate(text)


def _parse_pptx(file_bytes: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_bytes))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        parts = [f"--- Slide {i} ---"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
        slides.append("\n".join(parts))
    return _truncate("\n\n".join(slides))


def _parse_docx(file_bytes: bytes) -> str:
    from docx import Document

    doc = Document(io.BytesIO(file_bytes))
    text = "\n".join(p.text for p in doc.paragraphs)
    return _truncate(text)


def _parse_text(file_bytes: bytes) -> str:
    text = file_bytes.decode("utf-8", errors="replace")
    return _truncate(text)


def _truncate(text: str) -> str:
    if len(text) > MAX_CHARS:
        return text[:MAX_CHARS] + "\n\n[Document truncated to 50,000 characters]"
    return text
